# make_deck.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ---- Image files you said you have (same folder as this script) ----
IMG = {
    "cover":              "system-identification.png",
    "system_id":          "physical-modeling.png",
    "kalman_fusion":      "sensor-fusion.png",
    "motion_tracking":    "tracking.png",
    "regime_modes":       "transitioning.png",
    "hybrid_fusion":      "motion-categories.png",
}

# ---- Slide content (Title → Principle → Where → Impact) ----
SLIDES = [
    {
        "kind": "cover",
        "title": "Applied ML for Mobility & Device Location",
        "subtitle": "Signal Processing • Sensor Fusion • Real-Time Inference",
        "image": IMG["cover"],
    },
    {
        "title": "System Identification for Mobility & Physical Signals",
        "principle": "Understand the physics and signal behavior before modeling — frequency content, dynamics, noise.",
        "where": [
            "Volvo — FFT-based comfort signatures from acceleration data",
            "Chalmers — wave-based device modeling & regression on large datasets",
        ],
        "impact": [
            "Comfort inference as accurate as expert drivers",
            "Enabled automated gearbox tuning & HIL deployment",
        ],
        "image": IMG["system_id"],
    },
    {
        "title": "Kalman-Style Fusion & Confidence Tracking",
        "principle": "Predict → measure → correct + track uncertainty. Confidence is a first-class output.",
        "where": [
            "Volvo — uncertainty-triggered retraining in real-time vehicle tests",
            "Dialysis monitor (soft-sensor) — feed-forward thermal model corrected by temperature anchors",
        ],
        "impact": [
            "Stable real-time inference in noisy environments",
            "Replaced costly hardware with inference + sensors (analogy: inertial + GPS fusion)",
        ],
        "image": IMG["kalman_fusion"],
    },
    {
        "title": "Motion Detection, Trajectory Tracking & Mode Classification",
        "principle": "Sliding windows → features → event detection → online adaptation.",
        "where": [
            "Volvo — gear-shift event detection in vehicle signals",
            "SiB Solutions — object movement detection & temporal feature refinement",
        ],
        "impact": [
            "Continuously improving comfort model via driver feedback",
            "Lifted worst-case classification from ~34% to ~74% in a production product",
        ],
        "image": IMG["motion_tracking"],
    },
    {
        "title": "Regime Clustering & Operating Mode Segmentation",
        "principle": "Systems operate in modes — detect regimes and specialize models instead of forcing one global one.",
        "where": [
            "SiB Solutions — separate hard vs easy object-cases, targeted training",
            "Ericsson — mode-segmented network power analysis",
        ],
        "impact": [
            "Robust ML under distribution shifts",
            "Identified competitor performance inflection points in live networks",
        ],
        "image": IMG["regime_modes"],
    },
    {
        "title": "Hybrid GNSS + IMU + Opportunistic Anchors (Soft-Sensor Analogy)",
        "principle": "Dead-reckon when needed, reset drift when anchors appear — fuse model + measurements.",
        "where": [
            "Dialysis soft-sensor — thermodynamic observer + PT100 sensors as anchors",
            "Volvo feed-forward control — prediction loop corrected by real-time feedback",
        ],
        "impact": [
            "Replaced complex hardware with model-driven estimation",
            "Low-latency stability on physical systems with drift correction",
        ],
        "image": IMG["hybrid_fusion"],
    },
]

# ---- Helpers ----
def add_title_only(prs, title_text):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    return slide

def add_cover(prs, title, subtitle, image_path):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    # Title box
    left = Inches(0.7); top = Inches(0.6); width = Inches(10); height = Inches(1.2)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 20)

    # Subtitle
    st = tf.add_paragraph()
    st.text = subtitle
    st.level = 1
    st.font.size = Pt(20)
    st.font.color.rgb = RGBColor(70, 70, 70)
    st.space_before = Pt(6)

    # Image
    img_left = Inches(1.0)
    img_top  = Inches(2.0)
    img_width = Inches(10)
    slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

    return slide

def add_content_slide(prs, data):
    slide = add_title_only(prs, data["title"])

    # Left column: bullets
    left = Inches(0.7); top = Inches(1.6); width = Inches(6.3); height = Inches(4.8)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.clear()

    # Principle
    p = tf.paragraphs[0]
    p.text = "Principle: " + data["principle"]
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)

    # Where
    p = tf.add_paragraph()
    p.text = "Where I applied it:"
    p.font.size = Pt(16)
    p.font.bold = True
    for item in data["where"]:
        q = tf.add_paragraph()
        q.text = "• " + item
        q.level = 1
        q.font.size = Pt(16)

    # Impact
    p = tf.add_paragraph()
    p.text = "Impact:"
    p.font.size = Pt(16)
    p.font.bold = True
    for item in data["impact"]:
        q = tf.add_paragraph()
        q.text = "• " + item
        q.level = 1
        q.font.size = Pt(16)

    # Right column: image
    img_left = Inches(7.2); img_top = Inches(1.6); img_width = Inches(6.5)
    try:
        slide.shapes.add_picture(data["image"], img_left, img_top, width=img_width)
    except Exception:
        # If image missing, add a placeholder rectangle
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, img_left, img_top, img_width, Inches(4.5))
        rect.text = "Image not found: {}".format(data["image"])
        rect.text_frame.paragraphs[0].font.size = Pt(14)
        rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def build_deck(path="ML_signal_fusion_pitch.pptx"):
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cover
    add_cover(prs, SLIDES[0]["title"], SLIDES[0]["subtitle"], SLIDES[0]["image"])

    # Content slides
    for s in SLIDES[1:]:
        add_content_slide(prs, s)

    prs.save(path)
    print(f"Saved: {path}")

if __name__ == "__main__":
    build_deck()
